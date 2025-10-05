# -*- coding: utf-8 -*-
from pathlib import Path
import csv, json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent
EXP  = ROOT/"exports"

# ---- данные
analysis = (EXP/"analysis_Q4_2025.md").read_text(encoding="utf-8")
plan_rows = list(csv.DictReader((EXP/"plan_Q4_2025_justified.csv").open(encoding="utf-8")))
examples  = json.loads((EXP/"examples_ready_CLEAN.json").read_text(encoding="utf-8"))
targeting = json.loads((EXP/"targeting_recommendations.json").read_text(encoding="utf-8"))
ar_json   = (EXP/"visuals"/"ar_prompt.json").read_text(encoding="utf-8")
poster_png= (EXP/"visuals"/"poster_Immunocomplex.png")
by_week   = {}
for r in plan_rows:
    by_week.setdefault(r["week"], []).append(r)

# ---- стили
COL_BG, COL_ACC, COL_LIGHT, COL_DK, COL_TXT = RGBColor(10,24,48), RGBColor(0,163,163), RGBColor(250,250,252), RGBColor(40,48,60), RGBColor(28,28,28)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s, c): r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background()
def leftbar(s): b=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(0.35),prs.slide_height); b.fill.solid(); b.fill.fore_color.rgb=COL_ACC; b.line.fill.background()
def header(s, t):
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(12.3), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = COL_BG; band.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(0.58), Inches(11.7), Inches(0.9)).text_frame
    p = tb.paragraphs[0]; p.text = t
    for r in p.runs: r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=RGBColor(255,255,255)

def paragraph(tf, txt, size=14, color=COL_DK, bold=False, space=6):
    p = tf.add_paragraph() if tf.text != "" else tf.paragraphs[0]
    p.text = txt
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    p.space_after = Pt(space)
    return p

# ---- титул
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_BG)
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.4), Inches(2.1)).text_frame
tb.text = "Biotact IE-Agent — Q4 2025 (Prototype)"
for r in tb.paragraphs[0].runs: r.font.size=Pt(44); r.font.bold=True; r.font.color.rgb=RGBColor(255,255,255)
p2 = paragraph(tb, "12 недель • Обоснования • Примеры • AR • Таргетинг • Экспорт CSV/JSON", size=20, color=RGBColor(220,235,235))
# дисклеймер
td = s.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(11.4), Inches(1.2)).text_frame; td.text=""
paragraph(td,"Дисклеймер: материалы носят информационный характер и содержат формулировки «поддерживает/способствует». Не являются медицинскими рекомендациями.", size=12, color=RGBColor(220,235,235))
# ---- анализ (разбиваем на 2 слайда при длинном тексте)
chunks = re.split(r"\n\s*\n", analysis.strip())
for i,chunk in enumerate(chunks[:2], start=1):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s, f"Аналитический обзор Q4-2025 ({i}/2)")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.5)).text_frame; box.text=""
    for para in chunk.split("\n"):
        paragraph(box, para.strip(), size=16, color=COL_TXT, space=6)

# ---- выборка плана (таблица)
def short(d):
    m = re.match(r"(\d{4})-(\d{2})-(\d{2})", d or ""); return f"{m.group(3)}.{m.group(2)}" if m else (d or "")

s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s,"План контента — выборка (16 строк)")
rows = min(16, len(plan_rows)); cols = 6
tbl = s.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.7), Inches(12.4), Inches(5.6)).table
# ширины
tbl.columns[0].width = Inches(1.0)  # Week
tbl.columns[1].width = Inches(1.8)  # Dates
tbl.columns[2].width = Inches(1.8)  # Channel
tbl.columns[3].width = Inches(2.4)  # Product
tbl.columns[4].width = Inches(4.0)  # Topic
tbl.columns[5].width = Inches(1.4)  # Goal
heads = ["Week","Dates","Channel","Product","Topic","Goal"]
for j,h in enumerate(heads):
    cell = tbl.cell(0,j); cell.text = h
    for run in cell.text_frame.paragraphs[0].runs: run.font.size=Pt(12); run.font.bold=True
for i,r in enumerate(plan_rows[:rows], start=1):
    tbl.cell(i,0).text = r.get("week","")
    dates = f"{short(r.get('start',''))}–{short(r.get('end',''))}"
    tbl.cell(i,1).text = dates
    tbl.cell(i,2).text = r.get("channel","")
    tbl.cell(i,3).text = r.get("product","")
    topic = r.get("topic","")
    if len(topic) > 90: topic = topic[:87] + "…"
    tbl.cell(i,4).text = topic
    tbl.cell(i,5).text = r.get("goal","")
    for j in range(cols):
        for par in tbl.cell(i,j).text_frame.paragraphs:
            for run in par.runs: run.font.size=Pt(12)

# ---- недельные слайды (≤5 пунктов + rationale)
slide_no = 5
weeks_sorted = sorted(by_week.items(), key=lambda kv: int(kv[0]) if str(kv[0]).isdigit() else 999)
for week, items in weeks_sorted[:8]:   # показываем 8 недель как образец
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s, f"Неделя {week}: обоснование и задачи")
    # левая колонка — rationale
    left = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(6.1), Inches(5.6)).text_frame; left.text=""
    paragraph(left, "Обоснование", size=18, bold=True)
    rationale = (items[0].get("rationale","") or "—").strip()
    # режем до 5–6 предложений
    sentences = re.split(r"(?<=[.!?])\s+", rationale); rationale = " ".join(sentences[:6])
    paragraph(left, rationale, size=14, color=COL_TXT)
    # правая колонка — активности
    right = s.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(6.1), Inches(5.6)).text_frame; right.text=""
    paragraph(right, "Активности недели", size=18, bold=True)
    for it in items[:5]:
        line = f"• {it['channel']}: {it['topic']}  ({it['product']} → {it['format']}/{it['goal']})"
        paragraph(right, line, size=14, color=COL_TXT)
    slide_no += 1

# ---- примеры (очищенные)
for ex in examples:
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s, f"Примеры — {ex.get('sku','')}")
    left = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(6.1), Inches(5.6)).text_frame; left.text=""
    paragraph(left,"Instagram", size=18, bold=True)
    paragraph(left, ex.get("instagram",""), size=14, color=COL_TXT)
    right = s.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(6.1), Inches(5.6)).text_frame; right.text=""
    paragraph(right,"Email", size=18, bold=True);   paragraph(right, ex.get("email",""), size=14, color=COL_TXT)
    paragraph(right,"Podcast", size=18, bold=True); paragraph(right, ex.get("podcast",""), size=14, color=COL_TXT)

# ---- таргетинг
if targeting:
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s,"Таргетинг / Бюджет / Метрики")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.6)).text_frame; box.text=""
    for item in targeting:
        sku = item.get("sku",""); rec=item.get("recommendations",{}); aud=rec.get("audience",{}); met=rec.get("metrics",{})
        line = f"• {sku}: age {aud.get('age','')}, gender {aud.get('gender','')}, interests {', '.join(aud.get('interests',[]))}, geo {aud.get('geo','')}; budget ≈ {rec.get('budget_eur',0)}€; KPI: ER≥{met.get('ER_min_%','?')}%, CTR≥{met.get('CTR_min_%','?')}%, Conv≥{met.get('Conv_min_%','?')}%, Podcast≥{met.get('Podcast_watch_%','?')}%"
        paragraph(box, line, size=14, color=COL_TXT)

# ---- визуал/AR
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s,"Визуал / AR")
left = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(6.1), Inches(5.6)).text_frame; left.text=""
paragraph(left, "AR-промпт (суть)", size=18, bold=True)
short_ar = re.sub(r"\s+"," ", ar_json)[:450] + ("…" if len(ar_json)>450 else "")
paragraph(left, short_ar, size=14, color=COL_TXT)
# превью PNG (если есть)
if poster_png.exists():
    s.shapes.add_picture(str(poster_png), Inches(7.1), Inches(1.8), height=Inches(4.8))

# ---- интеграции/масштабируемость
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s,COL_LIGHT); leftbar(s); header(s,"Интеграция и масштабируемость")
box = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.6)).text_frame; box.text=""
for line in [
    "Экспорт CSV/JSON готов; Notion/Sheets — stubs (заменяются на реальные API).",
    "Интерактивный CLI: генерация анализа, плана, примеров, таргетинга, визуала.",
    "Модули по функциям (LLM с офлайн-фоллбэком).",
    "Roadmap: Telegram-бот, реальные интеграции, кеш LLM, AR-шаблоны."
]:
    paragraph(box, "• "+line, size=14, color=COL_TXT)

OUT = ROOT/"Biotact_Q4_2025_IE_AGENT_TWEAKED.pptx"
prs.save(OUT)
print("Saved:", OUT)
