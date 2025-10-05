# -*- coding: utf-8 -*-
# patch: увеличиваем шрифты и делаем короткий формат дат в таблице
from pathlib import Path
import csv, json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent
prs  = Presentation("Biotact_Q4_2025_IE_AGENT_FINAL.pptx")

# Найдём слайд с таблицей (заголовок "План контента — выборка")
for s in prs.slides:
    if any("План контента — выборка" in (sh.text if hasattr(sh,"text") else "") for sh in s.shapes):
        # первая таблица на слайде
        tbl = next(sh.table for sh in s.shapes if hasattr(sh, "table"))
        # укороченные даты и чуть крупнее шрифт
        def short(d):
            m = re.match(r"(\d{4})-(\d{2})-(\d{2})", d or "")
            return f"{m.group(3)}.{m.group(2)}" if m else (d or "")
        for i in range(1, len(tbl.rows)):
            # Week | Dates | Channel | Product | Topic | Goal
            tf = tbl.cell(i,1).text_frame
            tf.clear()
            p = tf.paragraphs[0]; p.text = f"{short(tbl.cell(i,1).text.split('–')[0])}–{short(tbl.cell(i,1).text.split('–')[-1])}"
            for r in p.runs: r.font.size = Pt(12)
            # уменьшить «Topic» до 90 символов
            t = tbl.cell(i,4).text
            if len(t) > 90: t = t[:87] + "…"
            tbl.cell(i,4).text = t

            # увеличить общий размер шрифта по строке
            for j in range(6):
                for par in tbl.cell(i,j).text_frame.paragraphs:
                    for run in par.runs:
                        run.font.size = Pt(12)

        # заголовки жирнее
        for j in range(6):
            for run in tbl.cell(0,j).text_frame.paragraphs[0].runs:
                run.font.size = Pt(12); run.font.bold = True
        break

prs.save("Biotact_Q4_2025_IE_AGENT_FINAL.pptx")
print("Reformatted table ✓")
