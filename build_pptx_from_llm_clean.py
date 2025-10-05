# -*- coding: utf-8 -*-
from pathlib import Path
import json, csv

ROOT = Path(__file__).parent
EXPORTS = ROOT / "exports"
EX_CLEAN = EXPORTS / "examples_llm_clean"
EX_RAW = EXPORTS / "examples_llm"
EX_SRC = EX_CLEAN if EX_CLEAN.exists() else EX_RAW
PLAN_FULL = EXPORTS / "plan_Q4_2025_full.csv"
PLAN = EXPORTS / "plan_Q4_2025.csv"
OUT = ROOT / "Biotact_Q4_2025_LLM_Styled_CLEAN.pptx"

insights = json.loads((EXPORTS / "insights.json").read_text(encoding="utf-8"))

plan_rows = []
plan_path = PLAN_FULL if PLAN_FULL.exists() else PLAN
with plan_path.open(encoding="utf-8") as f:
    rdr = csv.DictReader(f)
    plan_rows = list(rdr)
subset = plan_rows[:20]

examples = []
for p in sorted(EX_SRC.glob("example*_*.instagram.txt")):
    base = p.with_suffix("")
    sku = p.name.split("_", 1)[1].split(".")[0]
    ig = p.read_text(encoding="utf-8").strip()
    pod = base.with_suffix(".podcast.txt").read_text(encoding="utf-8").strip() if base.with_suffix(".podcast.txt").exists() else ""
    em = base.with_suffix(".email.txt").read_text(encoding="utf-8").strip() if base.with_suffix(".email.txt").exists() else ""
    arj = base.with_suffix(".ar.json").read_text(encoding="utf-8").strip() if base.with_suffix(".ar.json").exists() else ""
    examples.append({"sku": sku, "ig": ig, "pod": pod, "em": em, "ar": arj})

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

COL_BG = RGBColor(10, 24, 48)
COL_ACC = RGBColor(0, 163, 163)
COL_LIGHT = RGBColor(250,250,252)
COL_TEXT = RGBColor(255,255,255)
COL_DKTXT = RGBColor(40,48,60)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()

def leftbar(slide):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = COL_ACC; b.line.fill.background()

def header(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = COL_BG; band.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.9), Inches(0.9)).text_frame
    p = tb.paragraphs[0]; p.text = title
    for r in p.runs: r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=COL_TEXT

def footer(slide, n):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4)).text_frame
    tb.clear(); p = tb.paragraphs[0]; p.text = f"Generated • Slide {n}"
    for r in p.runs: r.font.size=Pt(10); r.font.color.rgb=COL_DKTXT

# Title
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_BG)
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(2.0)).text_frame
p = tb.paragraphs[0]; p.text = "Biotact Q4-2025 Marketing Plan + IE-Agent (LLM Clean)"
for r in p.runs: r.font.size=Pt(44); r.font.bold=True; r.font.color.rgb=COL_TEXT
p2 = tb.add_paragraph(); p2.text = "12-week plan • LLM examples (clean) • AR prompts"
for r in p2.runs: r.font.size=Pt(20); r.font.color.rgb=COL_TEXT
footer(s, 1)

# Insights
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s, "Seasonal Insights (Q4-2025)")
body = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.8), Inches(4.8)).text_frame
for i,it in enumerate(insights[:8]):
    p = body.add_paragraph() if i else body.paragraphs[0]
    p.text = f"• {it['sku']}: {it['insight']}"
    for r in p.runs: r.font.size=Pt(16); r.font.color.rgb=COL_DKTXT
footer(s, 2)

# Plan (sample)
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s, "Content Plan — sample (first 20 rows)")
rows = 1 + len(subset); cols = 5
table = s.shapes.add_table(rows, cols, Inches(0.8), Inches(2.0), Inches(12.0), Inches(4.8)).table
heads = ["Week","Dates","Channel","Product","Topic"]
for j,h in enumerate(heads):
    cell = table.cell(0,j); cell.text = h
    for r in cell.text_frame.paragraphs[0].runs: r.font.bold=True; r.font.size=Pt(13)
for i,r in enumerate(subset, start=1):
    table.cell(i,0).text = r.get("week","")
    table.cell(i,1).text = f"{r.get('start','')}–{r.get('end','')}"
    table.cell(i,2).text = r.get("channel","")
    table.cell(i,3).text = r.get("product","")
    table.cell(i,4).text = r.get("topic","")
footer(s, 3)

# LLM example slides
for idx, ex in enumerate(examples, start=4):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s, f"LLM Examples — {ex['sku']}")
    left = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.75), Inches(4.8)).text_frame
    p = left.paragraphs[0]; p.text = "Instagram"; 
    for r in p.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DKTXT
    p2 = left.add_paragraph(); p2.text = (ex["ig"][:900] + " …") if len(ex["ig"])>900 else ex["ig"]
    for r in p2.runs: r.font.size=Pt(13); r.font.color.rgb=COL_DKTXT

    right = s.shapes.add_textbox(Inches(6.55), Inches(2.0), Inches(6.25), Inches(4.8)).text_frame
    blocks = []
    if ex["pod"]: blocks.append(("Podcast", ex["pod"]))
    if ex["em"]:  blocks.append(("Email", ex["em"]))
    if ex["ar"]:  blocks.append(("AR JSON", ex["ar"]))
    if not blocks: blocks = [("Info","—")]
    for bi,(title,txt) in enumerate(blocks):
        p = right.add_paragraph() if bi>0 else right.paragraphs[0]
        p.text = title
        for r in p.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DKTXT
        p2 = right.add_paragraph(); p2.text = (txt[:900] + " …") if len(txt)>900 else txt
        for r in p2.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DKTXT
    footer(s, idx)

prs.save(OUT)
print("Saved:", OUT)
