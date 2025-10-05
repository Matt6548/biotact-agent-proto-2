# -*- coding: utf-8 -*-
from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Pt

SRC = "Biotact_Q4_2025_IE_AGENT_FINAL.pptx"
DST = "Biotact_Q4_2025_IE_AGENT_TWEAKED.pptx"

prs = Presentation(SRC)

# --- 1) Таблица "План контента — выборка": укрупнить шрифт, обрезать Topic, даты сделать короче
def short_date(d):
    m = re.match(r"(\d{4})-(\d{2})-(\d{2})", d or "")
    return f"{m.group(3)}.{m.group(2)}" if m else (d or "")

for s in prs.slides:
    title_texts = [getattr(sh, "text", "") for sh in s.shapes if hasattr(sh, "text")]
    if any("План контента — выборка" in t for t in title_texts):
        tbl = next(sh.table for sh in s.shapes if hasattr(sh, "table"))
        # заголовки
        for j in range(len(tbl.columns)):
            for run in tbl.cell(0,j).text_frame.paragraphs[0].runs:
                run.font.size = Pt(12); run.font.bold = True
        # строки
        for i in range(1, len(tbl.rows)):
            # Dates
            dates_raw = tbl.cell(i,1).text
            parts = dates_raw.split("–")
            left = short_date(parts[0].strip()) if parts else dates_raw
            right = short_date(parts[-1].strip()) if len(parts)>1 else ""
            tbl.cell(i,1).text = f"{left}–{right}" if right else left
            # Topic (обрезаем до 90 симв.)
            t = tbl.cell(i,4).text
            if len(t) > 90: t = t[:87] + "…"
            tbl.cell(i,4).text = t
            # общий размер шрифта
            for j in range(len(tbl.columns)):
                for par in tbl.cell(i,j).text_frame.paragraphs:
                    for run in par.runs:
                        run.font.size = Pt(12)
        break

# --- 2) Правые колонки "Активности недели": ≤5 пунктов и 14 pt
for s in prs.slides:
    for sh in s.shapes:
        if hasattr(sh, "text_frame") and "Активности недели" in (sh.text_frame.text or ""):
            lines = [p.text for p in sh.text_frame.paragraphs]
            head = lines[0] if lines else "Активности недели"
            items = [l for l in lines[1:] if l.strip().startswith("•")][:5]
            tf = sh.text_frame
            tf.clear()
            p = tf.paragraphs[0]; p.text = head
            for r in p.runs: r.font.size = Pt(16); r.font.bold = True
            for it in items:
                pp = tf.add_paragraph(); pp.text = it
                for r in pp.runs: r.font.size = Pt(14)

prs.save(DST)
print(f"Saved tweaked PPTX -> {DST}")
