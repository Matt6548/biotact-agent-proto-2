# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt
prs = Presentation("Biotact_Q4_2025_IE_AGENT_FINAL.pptx")
for s in prs.slides:
    if any("Активности недели" in (sh.text if hasattr(sh,"text") else "") for sh in s.shapes):
        # все текстбоксы на правой колонке: укрупняем и обрезаем лишние пункты
        boxes = [sh for sh in s.shapes if hasattr(sh, "text_frame")]
        for bx in boxes:
            if "Активности недели" in bx.text:
                # оставляем заголовок + первые 5 bullets
                lines = [p.text for p in bx.text_frame.paragraphs]
                head = lines[0]
                items = [l for l in lines[1:] if l.strip().startswith("•")][:5]
                bx.text_frame.clear()
                p = bx.text_frame.paragraphs[0]; p.text = head
                for r in p.runs: r.font.size = Pt(16); r.font.bold = True
                for it in items:
                    pp = bx.text_frame.add_paragraph(); pp.text = it
                    for r in pp.runs: r.font.size = Pt(14)
    # конец цикла по слайдам
prs.save("Biotact_Q4_2025_IE_AGENT_FINAL.pptx")
print("Bullets resized ✓")
