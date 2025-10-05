# -*- coding: utf-8 -*-
from pathlib import Path
import csv, json, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent
EXP = ROOT/"exports"

# --- Загрузка данных ---
analysis = (EXP/"analysis_Q4_2025.md").read_text(encoding="utf-8") if (EXP/"analysis_Q4_2025.md").exists() else ""
def trim_words(t, n):
    w = t.split()
    return t if len(w)<=n else " ".join(w[:n])+" …"

# План с обоснованиями
plan_rows = []
with (EXP/"plan_Q4_2025_justified.csv").open(encoding="utf-8") as f:
    rdr = csv.DictReader(f)
    plan_rows = list(rdr)

# Сгруппируем по неделям и соберём брифы
by_week = {}
for r in plan_rows:
    by_week.setdefault(r["week"], []).append(r)

# Примеры
examples = json.loads((EXP/"examples_ready.json").read_text(encoding="utf-8")) if (EXP/"examples_ready.json").exists() else []

# Таргетинг
targeting = json.loads((EXP/"targeting_recommendations.json").read_text(encoding="utf-8")) if (EXP/"targeting_recommendations.json").exists() else []

# Визуалы
ar_json = (EXP/"visuals"/"ar_prompt.json").read_text(encoding="utf-8") if (EXP/"visuals"/"ar_prompt.json").exists() else "{}"
poster = None
for p in (EXP/"visuals").glob("poster_*.svg"):
    poster = p
    break

# --- Стили ---
COL_BG = RGBColor(10,24,48)
COL_ACC = RGBColor(0,163,163)
COL_LIGHT = RGBColor(250,250,252)
COL_TEXT = RGBColor(255,255,255)
COL_DK = RGBColor(40,48,60)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()

def leftbar(slide):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = COL_ACC; b.line.fill.background()

def header(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = COL_BG; band.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.9), Inches(0.9)).text_frame
    p = tb.paragraphs[0]; p.text = title
    for run in p.runs: run.font.size=Pt(30); run.font.bold=True; run.font.color.rgb=COL_TEXT

def footer(slide, n):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4)).text_frame
    tb.clear(); p = tb.paragraphs[0]; p.text = f"Biotact • Slide {n}"
    for run in p.runs: run.font.size=Pt(10); run.font.color.rgb=COL_DK

# --- 1. Титул ---
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_BG)
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.5), Inches(2.0)).text_frame
p = tb.paragraphs[0]; p.text = "Biotact Q4-2025 IE-Agent (Prototype)"
for r in p.runs: r.font.size=Pt(44); r.font.bold=True; r.font.color.rgb=COL_TEXT
p2 = tb.add_paragraph(); p2.text = "План 12 недель • Обоснования • Примеры • AR • Таргетинг"
for r in p2.runs: r.font.size=Pt(20); r.font.color.rgb=COL_TEXT
footer(s,1)

# --- 2. Анализ (300–500 слов) ---
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s,"Аналитический обзор Q4-2025")
txt = trim_words(analysis, 480)
body = s.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(12.0), Inches(5.4)).text_frame
for i,para in enumerate(txt.split("\n")):
    p = body.add_paragraph() if i else body.paragraphs[0]
    p.text = para.strip()
    for r in p.runs: r.font.size=Pt(14); r.font.color.rgb=COL_DK
footer(s,2)

# --- 3. Календарь: обзор (первые 16 строк) ---
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s,"План контента — выборка")
rows = min(16, len(plan_rows)); cols = 6
table = s.shapes.add_table(rows+1, cols, Inches(0.6), Inches(1.8), Inches(12.2), Inches(5.4)).table
heads = ["Week","Dates","Channel","Product","Topic","Goal"]
for j,h in enumerate(heads):
    cell = table.cell(0,j); cell.text = h
    for r in cell.text_frame.paragraphs[0].runs: r.font.bold=True; r.font.size=Pt(12)
for i,r in enumerate(plan_rows[:rows], start=1):
    table.cell(i,0).text = r.get("week","")
    table.cell(i,1).text = f"{r.get('start','')}–{r.get('end','')}"
    table.cell(i,2).text = r.get("channel","")
    table.cell(i,3).text = r.get("product","")
    table.cell(i,4).text = r.get("topic","")
    table.cell(i,5).text = r.get("goal","")
footer(s,3)

# --- 4+. Слайды по неделям (rationale + 4–6 строк недели) ---
slide_no = 4
for week, items in sorted(by_week.items(), key=lambda kv: int(kv[0]) if kv[0].isdigit() else 999):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s, f"Неделя {week}: обоснование и задачи")
    # rationale (берём первый у недели)
    rationale = items[0].get("rationale","").strip()
    left = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.3), Inches(5.4)).text_frame
    p = left.paragraphs[0]; p.text = "Обоснование"
    for r in p.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DK
    p2 = left.add_paragraph(); p2.text = trim_words(rationale, 140) if rationale else "—"
    for r in p2.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK

    # правый столбец — строки недели
    right = s.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(6.2), Inches(5.4)).text_frame
    q = right.paragraphs[0]; q.text = "Активности недели"
    for r in q.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DK
    for it in items[:6]:
        line = f"• {it['channel']}: {it['topic']}  ({it['product']} → {it['format']}/{it['goal']})"
        p = right.add_paragraph(); p.text = line
        for rr in p.runs: rr.font.size=Pt(12); rr.font.color.rgb=COL_DK
    footer(s, slide_no); slide_no += 1
    if slide_no > 12: break  # чтобы не разрасталось — показываем первые ~8 недель как пример

# --- Слайды примеров (IG/Email/Podcast + AR) ---
for ex in examples:
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s, f"Примеры — {ex.get('sku','')}")
    left = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.2), Inches(5.4)).text_frame
    p = left.paragraphs[0]; p.text = "Instagram"
    for r in p.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DK
    p2 = left.add_paragraph(); p2.text = trim_words(ex.get("instagram",""), 130)
    for r in p2.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK

    right = s.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(6.2), Inches(5.4)).text_frame
    for title, key, limit in [("Email","email",130),("Podcast","podcast",160),("AR JSON","ar",90)]:
        h = right.add_paragraph(); h.text = title
        for r in h.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DK
        t = right.add_paragraph(); t.text = trim_words(ex.get(key,""), limit)
        for r in t.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK
    footer(s, slide_no); slide_no += 1

# --- Таргетинг/метрики слайд ---
if targeting:
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s,"Таргетинг и метрики (3 SKU)")
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.2), Inches(5.4)).text_frame
    for i,item in enumerate(targeting):
        p = body.add_paragraph() if i else body.paragraphs[0]
        sku = item.get("sku","")
        rec = item.get("recommendations",{})
        aud = rec.get("audience",{})
        met = rec.get("metrics",{})
        p.text = f"• {sku}: age {aud.get('age','')}, gender {aud.get('gender','')}, interests {', '.join(aud.get('interests',[]))}, geo {aud.get('geo','')}; budget ≈ {rec.get('budget_eur',0)}€; KPI: ER≥{met.get('ER_min_%','?')}%, CTR≥{met.get('CTR_min_%','?')}%, Conv≥{met.get('Conv_min_%','?')}%, Podcast≥{met.get('Podcast_watch_%','?')}%"
        for r in p.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK
    footer(s, slide_no); slide_no += 1

# --- Визуал/AR слайд ---
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s,"Визуал / AR")
box = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.2), Inches(5.4)).text_frame
p = box.paragraphs[0]; p.text = "AR-промпт (JSON)"
for r in p.runs: r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=COL_DK
p2 = box.add_paragraph(); p2.text = trim_words(ar_json.replace("\n"," "), 120) if ar_json else "—"
for r in p2.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK
if poster:
    p3 = box.add_paragraph(); p3.text = f"Постер (SVG): {poster.name} — см. в папке exports/visuals/"
    for r in p3.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK

# --- Итоги/интеграции ---
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, COL_LIGHT); leftbar(s); header(s,"Интеграция и масштабирование")
body = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.2), Inches(5.4)).text_frame
for i, line in enumerate([
    "Экспорт CSV/JSON готов (Notion/Sheets — stubs, легко заменить на реальные API).",
    "Кросс-канальная логика: соцсети → сайт → email → подкаст; AR/QR поддерживают вовлечение.",
    "Масштабируемость: отдельные модули под анализ, план, тексты, визуал, таргетинг.",
    "Дальше: подключить реальные API (Notion/Sheets/Telegram), кеш LLM, AR-шаблоны."
]):
    p = body.add_paragraph() if i else body.paragraphs[0]
    p.text = "• " + line
    for r in p.runs: r.font.size=Pt(12); r.font.color.rgb=COL_DK
footer(s, slide_no)

OUT = ROOT/"Biotact_Q4_2025_IE_AGENT_FINAL.pptx"
prs.save(OUT)
print("Saved:", OUT)
